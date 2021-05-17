"""
For License information see the LICENSE file.

Authors: Johannes Leupold, Tobias Stöckert, Amos Treiber

"""
import csv
import email
import io
import json
import re
import xml.etree.ElementTree as ET
from abc import ABC, abstractmethod
from collections import namedtuple
from email.message import Message
from logging import getLogger
from os import path, walk
from re import match
from typing import Iterator, Union, Callable, Iterable, TextIO, List, Tuple, Dict
from urlextract import URLExtract
from urllib.parse import urlparse
from datetime import datetime
from pptx import Presentation
from tika import parser as pdfparser

from ..pipeline import Source, Filter
from ...api import InputDocument, QueryInputDocument
try:
    from xml.etree.cElementTree import XML
except ImportError:
    from xml.etree.ElementTree import XML
import zipfile

log = getLogger(__name__)

RelativeFile = namedtuple("RelativeFile", ["base", "relative"])
LoadedFile = namedtuple("LoadedFile", ["relative_name", "content", "payload"])
WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
PARA = WORD_NAMESPACE + 'p'
TEXT = WORD_NAMESPACE + 't'


def filter_url(url: str) -> str:
    """
    Filters occurrences of urls/hyperlinks to only contain the domain name as a keyword

    Parameters
    ----------
    url : str
        the string containing an url

    Returns
    ----------
    url : str
        the url string with any url occurrences being replaced with just their domain name
    """
    res = url.strip()
    extractor = URLExtract()
    if extractor.has_urls(url):
        url = extractor.find_urls(url)[0]
        if "www." in url:
            t = urlparse(url)
            if t.netloc != '':
                t = t.netloc
                res = t.split('.')[1]
            elif t.path != '':
                t = t.path
                res = t.split('.')[1]
            else:
                res = url.replace("www.", "")

        else:
            t = urlparse(url).path
            res = t.split('.')[0]

    return res.strip()


"""
Method that extracts text from MS XML Word document (.docx).
(Inspired by python-docx <https://github.com/mikemaccana/python-docx>)
"""
def get_docx_text(path):
    """
    Take the path of a docx file as argument, return the text in unicode.
    """
    document = zipfile.ZipFile(path)
    xml_content = document.read('word/document.xml')
    document.close()
    tree = XML(xml_content)

    paragraphs = []
    for paragraph in tree.getiterator(PARA):
        texts = [node.text
                 for node in paragraph.getiterator(TEXT)
                 if node.text]
        if texts:
            paragraphs.append(''.join(texts))

    return '\n\n'.join(paragraphs)


def time_convert(timestr: str, timeformat: str) -> datetime:
    """
    Converts the given timestr to a datetime.datetime object according to the given timeformat
    """
    return datetime.strptime(timestr, timeformat)


class DirectoryEnumerator(Source[RelativeFile]):
    """
    A `Source` enumerating the files inside a base directory recursively and returning them as a named tuple
    `RelativeFile(base, relative)`, where `base` is the base directory (for later reference), and `relative` is
    the path of the file relative to the base directory.

    Consider for example the following directory hierarchy:

    /home/
        user/
            files/
                file1.txt
                file2.txt
                sub_directory/
                    file3.txt

    Then `DirectoryEnumerator('/home/user/files').elements()` would yield
    `RelativeFile('/home/user/files', 'file1.txt'), RelativeFile('/home/user/files', 'file2.txt'),
    RelativeFile('/home/user/files', 'sub_directory/file3.txt')`

    Parameters
    ----------
    base: str
        The base directory to enumerate,
        default: the working directory
    """

    def __init__(self, base: str = "."):
        self.__base: str = base

    def elements(self) -> Iterator[RelativeFile]:
        for dirname, _, files in walk(self.__base):
            for f in files:
                relative_dir = dirname[len(self.__base):].lstrip(path.sep)
                relative_file = path.join(relative_dir, f)
                yield RelativeFile(base=self.__base, relative=relative_file.replace(path.sep, "/"))


class RelativePrefixFilter(Filter[RelativeFile, RelativeFile]):
    """
    A `Filter` for `RelativeFile`, filtering for specific prefixes on the relative path component. It can be
    instantiated with a single string or any iterable of strings and will then yield all `RelativeFile` from the input
    iterator that have a prefix which matches that exact string or matches any string in the iterable.

    Parameters
    ----------
    match_prefix: Union[str, Iterable[str]]
        all prefixes to keep
    """

    def __init__(self, match_prefix: Union[str, Iterable[str]]):
        self.__match_predicate: Callable[[str], bool] = RelativePrefixFilter.__create_match_predicate(match_prefix)

    def filter(self, source: Iterator[RelativeFile]) -> Iterator[RelativeFile]:
        for f in source:
            if self.__match_predicate(f.relative):
                yield f

    @staticmethod
    def __create_match_predicate(match_prefix: Union[str, Iterable[str]]) -> Callable[[str], bool]:
        if isinstance(match_prefix, str):
            str_prefix: str = match_prefix
            return lambda rel: rel.startswith(str_prefix)
        else:
            return lambda rel: any(rel.startswith(prefix) for prefix in match_prefix)


class RelativeContainsFilter(Filter[RelativeFile, RelativeFile]):
    """
    A `Filter` for `RelativeFile`, filtering for specific occurrences on the relative path component. It can be
    instantiated with a single string or any iterable of strings and will then yield all `RelativeFile` from the input
    iterator that match that exact string or matches any string in the iterable.

    Parameters
    ----------
    matches: Union[str, Iterable[str]]
        all prefixes to keep
    """

    def __init__(self, matches: Union[str, Iterable[str]]):
        self.__match_predicate: Callable[[str], bool] = RelativeContainsFilter.__create_match_predicate(matches)

    def filter(self, source: Iterator[RelativeFile]) -> Iterator[RelativeFile]:
        for f in source:
            if self.__match_predicate(f.relative):
                yield f

    @staticmethod
    def __create_match_predicate(matches: Union[str, Iterable[str]]) -> Callable[[str], bool]:
        if isinstance(matches, str):
            str_contain: str = matches
            return lambda rel: str_contain in rel
        else:
            return lambda rel: any(m in rel for m in matches)


class FileParser(ABC):
    """
    A parser for parsing text resources to strings.
    """

    __doc_counter: int

    def __init__(self):
        self.__doc_counter = 0

    @abstractmethod
    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        """
        Takes a text resource and returns the content as Tuple[new_file_name="", content, payload=""] according to this
        parser. new_file_name and payload can be empty and are optional to manually alter the document name or to pass a
        payload associated with the content, like e.g., a user id.

        Parameters
        ----------
        f: TextIO
            the input resource
        """
        raise NotImplementedError

    def __call__(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        return self.parse(f)

    def incr_doc_count(self) -> None:
        self.__doc_counter += 1

    def doc_count(self) -> int:
        return self.__doc_counter


class PlainFileParser(FileParser):
    """
    A parser for plain text files that simply concatenates the lines of the file, stripping line breaks and adding
    spaces between the lines.
    """

    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        yield ("", " ".join([line for line in f.readlines()]), "")


class StructuredFileParser(FileParser, ABC):
    """
    A parser for structured files that takes specified field ids to produce document inputs for each line  in a structured
    document according to the specified attribute/field ids
    """
    _id_attribute: Union[str, int, None]
    _content_attribute: Union[str, int, List[int], List[str]]
    _payload_attribute: Union[str, int, None]

    def __init__(self, id_attribute: Union[str, int, None], content_attribute: Union[str, int, List[int]],
                 payload_attribute: Union[str, int, None] = None):
        super(StructuredFileParser, self).__init__()
        self._id_attribute = id_attribute
        self._content_attribute = content_attribute
        self._payload_attribute = payload_attribute


class XMLParser(StructuredFileParser):
    """
    A parser for xml files that takes specified field attributes to produce document inputs for each entry according to
    specified xml attribute names
    """

    _doc_attribute: str

    def __init__(self, id_attribute_name: str = 'id', content_attribute_name: Union[str, List[str]] = 'content',
                 payload_attribute_name: str = None, doc_attribute_name: str = 'Document'):
        super(XMLParser, self).__init__(id_attribute_name, content_attribute_name, payload_attribute_name)
        self._doc_attribute = doc_attribute_name

    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        tree = ET.parse(f)
        root = tree.getroot()
        i: int = 0
        for doc in root.findall('.//' + self._doc_attribute):  # './/' to find in all descendants
            id = ""
            if self._id_attribute is not None:
                find = doc.find(self._id_attribute)
                if find is not None:
                    id = find.text
            else:
                id = f"{self.doc_count()}_{i}"
            payload = ""
            if self._payload_attribute is not None:
                find = doc.find(self._payload_attribute)
                if find is not None:
                    payload = find.text

            content = ""
            if isinstance(self._content_attribute, str):
                find = doc.find(self._content_attribute)
                if find is not None:
                    payload = find.text
                content = find.text
            else:
                for c_attr in self._content_attribute:
                    contents = doc.findall('.//' + c_attr)
                    for c in contents:
                        content += " " + c.text

            i += 1
            yield (id, content, payload)


class JSONParser(StructuredFileParser):
    """
    A parser for json files that takes specified field ids to produce document inputs for each json line according to
    specified json attribute names
    """

    def __init__(self, id_attribute_name: str = 'id', content_attribute_name: str = 'content',
                 payload_attribute_name: str = None):
        super(JSONParser, self).__init__(id_attribute_name, content_attribute_name, payload_attribute_name)

    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        for line in f.readlines():
            doc = json.loads(line)
            payload = ""
            if self._payload_attribute is not None:
                payload = doc[self._payload_attribute].strip()
            yield (doc[self._id_attribute].strip(), doc[self._content_attribute], payload)


class CsvParser(StructuredFileParser):
    """
    A parser for csv files that takes specified field ids to produce document inputs for each csv line according to
    specified csv field ids
    """
    _delimiter: str

    def __init__(self, id_attribute_pos: Union[None, int] = None, content_attribute_pos: Union[int, List[int]] = 1,
                 payload_attribute_pos: Union[None, int] = 2, delimiter: str = '\t'):
        super(CsvParser, self).__init__(id_attribute_pos, content_attribute_pos, payload_attribute_pos)
        self._delimiter = delimiter

    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        csv_reader = csv.reader(f, delimiter=self._delimiter)
        next(csv_reader, None)  # skip first line

        id_attr = 0 if self._id_attribute is None else self._id_attribute
        payload_attr = 0 if self._payload_attribute is None else self._payload_attribute
        if isinstance(self._content_attribute, int):
            content_attr = self._content_attribute
        else:
            content_attr = min(self._content_attribute)  # We ignore content entries only if there are none

        for i, line in enumerate(csv_reader):
            if len(line) < max(id_attr, content_attr, payload_attr) + 1:
                log.warning(f"Encountered incorrect line {line} at position {self.doc_count()}_{i}. Skipping that")
                continue
            if self._id_attribute is None:  # use incremental int file ids FILENO_LINENO
                doc_id = f"{self.doc_count()}_{i}"
            else:
                doc_id = line[self._id_attribute].strip()
            payload = ""
            if self._payload_attribute is not None:
                payload = line[self._payload_attribute].strip()

            content = ""
            if isinstance(self._content_attribute, int):
                content = line[self._content_attribute]
            else:
                for c_attr in self._content_attribute:
                    if len(line) >= c_attr + 1:
                        content += " " + line[c_attr]

            yield (doc_id, content, payload)

        self.incr_doc_count()


class PubMedParser(CsvParser):
    """
    A parser for the PubMed query log. Since there are no keywords available,
    it outputs the recorded selectivities as keywords
    """

    __action_attribute_pos: int

    def __init__(self):
        super(PubMedParser, self).__init__(id_attribute_pos=None, content_attribute_pos=4, payload_attribute_pos=0,
                                           delimiter='\t')
        self.__action_attribute_pos = 2

    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        csv_reader = csv.reader(f, delimiter=self._delimiter)
        next(csv_reader, None)  # skip first line

        i = 0
        for line in csv_reader:
            if len(line) < 5:
                log.warning(f"Encountered incorrect line {line} at position {self.doc_count()}_{i}. Skipping that")
                continue
            if line[self.__action_attribute_pos] != "query":  # skip abstract and fulltext requests
                continue

            doc_id = f"{self.doc_count()}_{i}"

            payload = line[self._payload_attribute].strip()

            content = line[self._content_attribute].strip()

            yield (doc_id, content, payload)
            i += 1

        self.incr_doc_count()


class PocketDataParser(CsvParser):
    """
    A parser for the PocketData query log. Since there is no dataset available,
    it outputs the recorded selectivities as keywords
    """
    def __init__(self):
        super(PocketDataParser, self).__init__(id_attribute_pos=None, content_attribute_pos=8,
                                               payload_attribute_pos=None, delimiter='\t')

    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        csv_reader = csv.reader(f, delimiter=self._delimiter)
        next(csv_reader, None)  # skip first line

        i = 0
        for line in csv_reader:
            if len(line) < 9:
                log.warning(f"Encountered incorrect line {line} at position {self.doc_count()}_{i}. Skipping that")
                continue

            try:
                query: Dict[str, Union[int, str]] = eval(line[self._content_attribute])
            except SyntaxError:
                log.warning(f"Encountered incorrect line {line} at position {self.doc_count()}_{i}. Skipping that")
                continue

            if "Action" not in query:
                continue

            if query["Action"] != 'SELECT':
                continue  # skip non-select queries.

            selstr = "Rows returned"
            if selstr in query:
                resstr = "Results"
                if resstr in query:
                    if "FROM" in query[resstr]:
                        dbstr = query[resstr]
                        dbstr = dbstr[dbstr.find("FROM"):].split(" ")[1]
                        yield (f"{self.doc_count()}_{i}", f"{query[selstr]}", f"{self.doc_count()}_{dbstr}")
                        # since we have a file per user, use doc_count(). See separate DBs as different instances
                        # => Use doc_count()_DB as userid
                        i += 1

        self.incr_doc_count()


class AOLParser(CsvParser):
    """
    A parser for csv files of the AOL query log
    """

    __prev_content: str
    __prev_time: datetime

    def __init__(self):
        super(AOLParser, self).__init__(None, content_attribute_pos=1, payload_attribute_pos=0, delimiter='\t')
        self.__prev_content = ""
        self.__prev_datetime = time_convert("2000-01-01 00:00:00", timeformat="%Y-%m-%d %X")

    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        csv_reader = csv.reader(f, delimiter=self._delimiter)
        next(csv_reader, None)  # skip first line

        i: int = 0
        for line in csv_reader:
            content = line[self._content_attribute]
            time = time_convert(line[2], timeformat="%Y-%m-%d %X")
            if content == self.__prev_content and line[3] != "" and not (time - self.__prev_datetime).days >= 1:
                # the line represents a click and not a query
                continue

            doc_id = f"{self.doc_count()}_{i}"
            payload = line[self._payload_attribute].strip()

            i += 1
            self.__prev_content = content
            self.__prev_datetime = time
            yield (doc_id, filter_url(content), payload)

        self.incr_doc_count()


class EMailParser(FileParser):
    """
    A parser for e-mail files aiming at excluding e-mail headers, e-mail addresses and placeholders inserted for
    images.
    """

    @staticmethod
    def _payload_to_string(payload: Union[bytes, str, None, List[Message]]) -> str:
        if isinstance(payload, bytes):
            try:
                return payload.decode("utf8")
            except UnicodeDecodeError:
                log.warning(f"Encountered invalid bytestr. This line will be skipped.")
                return ""
        elif isinstance(payload, str):
            return payload
        return ""

    @staticmethod
    def _filtered(line: str) -> bool:
        if "-----" in line:
            return True
        elif line == "":
            return True
        elif match(r"^.*<.*@.*>.*$", line) is not None:
            return True
        elif match(r"^.*@.* on [0-9]{2}/[0-9]{2}/[0-9]{4}.*$", line) is not None:
            return True
        elif line.startswith(("From:", "FROM:", "from:", "To:", "TO:", "to:", "Cc:", "CC:", "cc:",
                              "Subject:", "SUBJECT:", "subject:")):
            return True
        elif match(r"^<<[^>]+>>", line) is not None:
            return True
        elif "[IMAGE]" in line:
            return True
        return False

    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        mail = email.message_from_file(fp=f)

        payload: str = ""
        if mail.is_multipart():
            for part in mail.walk():
                content_type = part.get_content_type()
                content_disposition = str(part.get('Content-Disposition'))

                # skip any text/plain (txt) attachments
                if content_type == 'text/plain' and 'attachment' not in content_disposition:
                    payload = EMailParser._payload_to_string(part.get_payload(decode=True))  # decode
                    break
        # not multipart - i.e. plain text, no attachments
        else:
            payload = EMailParser._payload_to_string(mail.get_payload(decode=True))

        yield ("", " ".join([line for line in payload.split("\n") if not EMailParser._filtered(line.strip())]),
               "")


class GoogleLogParser(JSONParser):
    """
    A parser for google acitivites downloaded via google takeout. Aims to return the searched keywords for activities
    in drive and gmail
    """
    def __init__(self, id_attribute_name: str = 'time', content_attribute_name: str = 'title',
                 payload_attribute_name: str = None):
        super(GoogleLogParser, self).__init__(id_attribute_name, content_attribute_name, payload_attribute_name)

    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        doc = json.loads(f.read())
        for dictionary in doc:
            payload = "0"
            if self._payload_attribute is not None:
                payload = dictionary[self._payload_attribute].strip()
            if re.match(r'^Searched for ', dictionary[self._content_attribute]):
                yield (dictionary[self._id_attribute].strip(), re.sub(r'^Searched for ', '',
                                                                      dictionary[self._content_attribute]),
                       payload)


class GMailParser(EMailParser):
    """
    A parser for .mbox e-mail files aiming at excluding e-mail headers, e-mail addresses and placeholders inserted for
    images.
    """

    def parse(self, f: TextIO) -> Iterator[Tuple[str, str, str]]:
        m: str = ""
        i: int = 0
        for l in f:
            if re.match(r"^From .*@xxx.*$", l):
                """New email"""
                mail = email.message_from_string(m)
                payload: str = ""
                if mail.is_multipart():
                    for part in mail.walk():
                        content_type = part.get_content_type()
                        content_disposition = str(part.get('Content-Disposition'))

                        # skip any text/plain (txt) attachments
                        if content_type == 'text/plain' and 'attachment' not in content_disposition:
                            payload = GMailParser._payload_to_string(part.get_payload(decode=True))  # decode
                            break
                # not multipart - i.e. plain text, no attachments, keeping fingers crossed
                else:
                    payload = GMailParser._payload_to_string(mail.get_payload(decode=True))

                i += 1
                m = ""

                yield (str(i),
                       " ".join([line for line in payload.split("\n") if not GMailParser._filtered(line.strip())]), "")

            else:
                """Line belongs to prior email"""
                m += l


class FileLoader(Filter[RelativeFile, LoadedFile]):
    """
    A `Filter` loading the content of all RelativeFiles found in the source. It will by default use the
    `PlainFileParser` for loading the files, but this behavior can be changed by specifying the `parser` parameter.
    PDF, docx, and pptx files are converted to text before being returned. A flag specifies whether other files should
    be ignored. The files are returned as a named tuple `LoadedFile(relative, content)`, where
    `relative` is the identically named     component of the `RelativeFile` (or, a new name returned by the fitler) and
    `content` is the content of the file as a string.

    Parameters
    ----------
    parser: Callable[[TextIO], Iterator[Tuple[str, str]]]
        The parser to use for loading files (new_filename, file_content). It may be a `FileParser`.
        default: `PlainFileParser()`
    parse_all: bool
        Whether unrecognized file types shall be parsed or shall be ignored
    """

    def __init__(self, parser: Callable[[TextIO], Iterator[Tuple[str, str, str]]] = PlainFileParser(),
                 parse_all: bool = True):
        self.__parser: Callable[[TextIO], Iterator[Tuple[str, str, str]]] = parser
        self.__parse_all: bool = parse_all

    def filter(self, source: Iterator[RelativeFile]) -> Iterator[LoadedFile]:
        for base, relative in source:
            txt = ""
            if re.match(r".*\.(docx)$", relative):
                txt = get_docx_text(path.join(base, relative))
            elif re.match(r".*\.(pptx)$", relative):
                txt = ""
                prs = Presentation(path.join(base, relative))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            txt = txt + shape.text + "\n"
            elif re.match(r".*\.(pdf)$", relative):
                raw = pdfparser.from_file(path.join(base, relative))
                txt = raw['content']
            elif re.match(r".*\.(txt)$", relative) or self.__parse_all:
                with open(path.join(base, relative), encoding="latin-1", errors="ignore") as f:
                    for new_filename, content, payload in self.__parser(f):
                        relative_name = relative + new_filename
                        yield LoadedFile(relative_name=relative_name, content=content, payload=payload)
            if txt != "":
                f = io.StringIO(txt)
                for new_filename, content, payload in self.__parser(f):
                    relative_name = relative + new_filename
                    yield LoadedFile(relative_name=relative_name, content=content, payload=payload)


class StripPrefix(Filter[LoadedFile, LoadedFile]):
    """
    A `Filter` capable of removing a common prefix from the `relative` component of `LoadedFile`. A file not having
    that prefix will be reyielded unchanged.

    Parameters
    ----------
    prefix: str
        The prefix to remove.
    """

    def __init__(self, prefix: str):
        self.__prefix: str = prefix

    def filter(self, source: Iterator[LoadedFile]) -> Iterator[LoadedFile]:
        for f in source:
            if f.relative_name.startswith(self.__prefix):
                yield f._replace(relative_name=f.relative_name[len(self.__prefix):])
            else:
                yield f


class FileToDocument(Filter[LoadedFile, InputDocument]):
    """
    A `Filter` converting `LoadedFile` to `InputDocument`. Ignores any payloads.
    """

    def filter(self, source: Iterator[LoadedFile]) -> Iterator[InputDocument]:
        for relative_name, content, _ in source:
            if len(content) == 0:
                log.warning(f"Encountered entry {relative_name} without any content. Skipping that...")
                continue
            yield InputDocument(relative_name, content)


class FileToQueryInputDocument(Filter[LoadedFile, QueryInputDocument]):
    """
    A `Filter` converting `LoadedFile` to `InputDocument`. Uses payload as user id.
    """

    def filter(self, source: Iterator[LoadedFile]) -> Iterator[QueryInputDocument]:
        for relative_name, content, payload in source:
            if len(content) == 0:
                log.warning(f"Encountered entry {relative_name} without any content. Skipping that...")
                continue
            yield QueryInputDocument(relative_name, content, payload)
